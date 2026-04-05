from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/seoki/Desktop/research")
OUTPUT_ROOT = ROOT / "outputs" / "presentation_deck_outline_61day_2026-03-13"
WORKBENCH = OUTPUT_ROOT / "workbenches" / "presentation_workbench_main_61day"
GENERATED = WORKBENCH / "generated_assets"

HOUSTON = OUTPUT_ROOT / "conference_print_assets_61day" / "figure1_holdout_compare_houston.png"
SCENARIO = OUTPUT_ROOT / "conference_print_assets_61day" / "figure2_scenario_aware_contour_compare.png"
REGIONAL = OUTPUT_ROOT / "conference_print_assets_61day" / "figure3_regional_failure_mode_schematic.png"
NOLA = GENERATED / "nola_holdout.png"
SEATTLE = GENERATED / "seattle_holdout.png"

OUT_8 = WORKBENCH / "presentation_8slide_production_draft_61day.pptx"
OUT_3 = WORKBENCH / "presentation_3minute_production_draft_61day.pptx"

BG = RGBColor(244, 241, 232)
INK = RGBColor(29, 43, 42)
SUB = RGBColor(90, 100, 95)
ACCENT = RGBColor(34, 92, 85)
ACCENT_2 = RGBColor(185, 136, 71)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(208, 203, 191)


def ensure_assets() -> None:
    missing = [p for p in [HOUSTON, SCENARIO, REGIONAL, NOLA, SEATTLE] if not p.exists()]
    if missing:
        raise FileNotFoundError(f"Missing assets: {missing}")


def base_presentation() -> tuple[Presentation, object]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_textbox(slide, x, y, w, h, text, *, size=20, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, bullets, *, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(7)
        r = p.runs[0]
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_header(slide, title, key):
    add_textbox(slide, 0.55, 0.3, 9.6, 0.42, title, size=28, bold=True)
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.72),
        Inches(0.28),
        Inches(1.9),
        Inches(0.42),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT
    tag.line.color.rgb = ACCENT
    tf = tag.text_frame
    tf.text = key
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.82), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_callout(slide, x, y, w, h, title, body, fill):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    r1 = p1.runs[0]
    r1.font.name = "Aptos"
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = body
    r2 = p2.runs[0]
    r2.font.name = "Aptos"
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = WHITE


def add_table(slide, x, y, w, h, data, *, col_widths=None, font_size=14):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(225, 232, 229) if r == 0 else WHITE
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = INK
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return table


def add_pipeline_box(slide, x, y, w, h, text, fill):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(14.5)
    r.font.bold = True
    r.font.color.rgb = WHITE


def build_8_slide(prs, blank):
    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "왜 공간 표현인가", "Problem")
    add_textbox(slide, 0.7, 1.2, 11.7, 0.6, "문제는 위험 계산 부족이 아니라 위험 공간 표현 부족이다.", size=25, bold=True, color=ACCENT)
    add_bullets(slide, 0.9, 2.05, 6.1, 3.1, [
        "기존 AIS 분석은 pairwise 숫자 판단에 강하다.",
        "실제 의사결정은 \"내 배 주변 어디가 더 위험한가\"라는 공간 질문에 가깝다.",
        "본 프로젝트는 pairwise risk를 spatial field로 바꾼다.",
    ], size=19)
    add_callout(slide, 7.6, 2.0, 4.7, 2.5, "Bridge Question", "Not just whether another vessel is risky, but which surrounding area becomes more risky for the own ship.", ACCENT)
    add_textbox(slide, 7.7, 5.0, 4.5, 0.8, "Pairwise risk -> own-ship-centric spatial field", size=18, bold=True, color=ACCENT_2)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "제안 구조", "Method")
    add_textbox(slide, 0.7, 1.15, 11.5, 0.55, "pairwise relation을 공간 representation으로 재구성한다.", size=24, bold=True, color=ACCENT)
    add_bullets(slide, 0.8, 2.0, 3.6, 2.5, ["trajectory reconstruction", "pairwise severity estimation", "heatmap + contour generation"], size=18)
    steps = ["Curate AIS", "Reconstruct Tracks", "Estimate Pairwise Risk", "Generate Heatmap & Contour"]
    for i, text in enumerate(steps):
        add_pipeline_box(slide, 4.7 + i * 2.0, 2.4, 1.75, 1.1, text, ACCENT if i % 2 == 0 else ACCENT_2)
    add_textbox(slide, 4.9, 4.1, 6.9, 1.0, "도메인 파이프라인 위에 설명 가능한 위험 표현을 고정하고, 모델은 severity layer에만 얹는다.", size=17, color=SUB)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "구현 파이프라인", "Pipeline")
    add_textbox(slide, 0.7, 1.15, 11.6, 0.55, "도메인 파이프라인 위에 모델이 올라간다.", size=24, bold=True, color=ACCENT)
    labels = ["focus", "tracks", "pairwise", "scenario shift", "threshold tuning"]
    xs = [0.8, 3.0, 5.2, 7.6, 10.0]
    for i, (x, label) in enumerate(zip(xs, labels)):
        add_pipeline_box(slide, x, 2.4, 2.0, 1.15, label, ACCENT if i % 2 == 0 else ACCENT_2)
    add_textbox(slide, 0.9, 4.25, 11.2, 0.9, "NOAA 수집부터 contour 비교까지 반복 가능한 구조를 만들었다.", size=18, color=SUB)
    add_textbox(slide, 0.9, 5.15, 11.4, 1.1, "focus -> tracks -> pairwise -> scenario shift -> threshold tuning", size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "검증 설계", "Validation")
    add_textbox(slide, 0.7, 1.1, 11.7, 0.55, "single split이 아니라 split + LOO + repeat를 함께 봤다.", size=24, bold=True, color=ACCENT)
    add_table(slide, 0.8, 2.0, 4.7, 2.25, [
        ["항목", "값"],
        ["기간", "61일"],
        ["해역", "Houston / NOLA / Seattle"],
        ["own_mmsi", "15"],
        ["pairwise rows", "67,892"],
    ], col_widths=[2.1, 2.5], font_size=14)
    add_callout(slide, 6.5, 2.0, 2.0, 1.1, "Benchmark", "timestamp split", ACCENT)
    add_callout(slide, 8.8, 2.0, 2.0, 1.1, "LOO", "own-ship holdout", ACCENT_2)
    add_callout(slide, 7.65, 3.45, 2.0, 1.1, "Repeat", "case repeat", ACCENT)
    add_textbox(slide, 6.1, 5.1, 5.8, 0.8, "benchmark + own-ship LOO + case repeat", size=19, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "모델 결과", "Results")
    add_textbox(slide, 0.7, 1.05, 11.7, 0.55, "주력 모델은 hgbt, 설명 가능한 비교군은 logreg다.", size=24, bold=True, color=ACCENT)
    add_table(slide, 0.75, 1.95, 8.1, 2.45, [
        ["Model", "Benchmark F1", "ECE", "LOO F1 mean"],
        ["hgbt", "0.9453", "0.0187", "0.9241"],
        ["logreg", "0.8750", "0.0644", "0.8953"],
    ], col_widths=[1.6, 2.0, 1.4, 2.2], font_size=15)
    add_callout(slide, 9.2, 2.0, 3.1, 1.15, "Primary", "hgbt", ACCENT)
    add_callout(slide, 9.2, 3.35, 3.1, 1.15, "Comparator", "logreg", ACCENT_2)
    add_bullets(slide, 0.95, 4.95, 10.8, 1.4, [
        "모델 layer에서는 hgbt가 가장 안정적이었다.",
        "설명 가능성은 logreg comparator로 보완했다.",
    ], size=17, color=SUB)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Threshold 결과", "Threshold")
    add_textbox(slide, 0.7, 1.05, 11.7, 0.55, "single best value 대신 shortlist가 더 정직했다.", size=24, bold=True, color=ACCENT)
    add_table(slide, 0.75, 1.95, 4.6, 2.3, [
        ["항목", "값"],
        ["majority profile", "s0p30_w0p55"],
        ["majority ratio", "0.1803"],
        ["mean top-k Jaccard", "0.1388"],
    ], col_widths=[2.2, 2.1], font_size=14)
    add_callout(slide, 5.8, 1.95, 2.1, 1.0, "default", "s0p30_w0p55", ACCENT)
    add_callout(slide, 8.05, 1.95, 2.1, 1.0, "balanced", "s0p30_w0p65", ACCENT_2)
    add_callout(slide, 10.3, 1.95, 2.1, 1.0, "tight", "s0p35_w0p65", ACCENT)
    slide.shapes.add_picture(str(SCENARIO), Inches(5.8), Inches(3.0), width=Inches(6.0))
    add_textbox(slide, 0.85, 4.7, 4.2, 1.0, "unstable -> shortlist 운영", size=21, bold=True, color=ACCENT_2)

    # 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Holdout Figure", "Visual")
    add_textbox(slide, 0.7, 1.0, 11.8, 0.55, "shortlist는 위험 방향보다 contour 민감도를 조정한다.", size=24, bold=True, color=ACCENT)
    slide.shapes.add_picture(str(HOUSTON), Inches(0.7), Inches(1.75), width=Inches(7.1))
    slide.shapes.add_picture(str(NOLA), Inches(8.15), Inches(1.95), width=Inches(2.2))
    slide.shapes.add_picture(str(SEATTLE), Inches(10.5), Inches(1.95), width=Inches(2.2))
    add_bullets(slide, 8.1, 4.65, 4.6, 1.5, [
        "default -> balanced: warning area 축소",
        "balanced -> tight: caution area 추가 축소",
        "dominant sector는 대체로 유지",
    ], size=15)

    # 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "기여와 한계", "Close")
    add_textbox(slide, 0.7, 1.0, 11.8, 0.55, "강점은 과장 없는 AIS-only decision-support framing이다.", size=24, bold=True, color=ACCENT)
    add_callout(slide, 0.8, 1.95, 5.7, 3.0, "기여", "spatial field\nscenario-aware contour\nregional failure-mode reporting", ACCENT)
    add_callout(slide, 6.8, 1.95, 5.7, 3.0, "비주장", "autonomy\nlegal safety boundary\nsingle optimum threshold", ACCENT_2)
    add_textbox(slide, 1.0, 5.45, 11.2, 0.8, "Conclusion: explainable spatial decision support", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(str(REGIONAL), Inches(3.75), Inches(6.0), width=Inches(5.7))


def build_3_slide_set(prs, blank):
    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "왜 공간 표현인가", "Problem")
    add_textbox(slide, 0.7, 1.25, 11.7, 0.55, "문제는 위험 계산 부족이 아니라 위험 공간 표현 부족이다.", size=25, bold=True, color=ACCENT)
    add_bullets(slide, 0.9, 2.15, 11.0, 2.6, [
        "기존 AIS 분석은 pairwise 숫자 판단에 강하다.",
        "실제 의사결정은 \"내 배 주변 어디가 더 위험한가\"라는 공간 질문에 가깝다.",
        "본 프로젝트는 pairwise risk를 own-ship-centric spatial field로 바꾼다.",
    ], size=20)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "제안 구조와 검증", "Method")
    add_textbox(slide, 0.7, 1.05, 11.8, 0.55, "도메인 파이프라인 위에서 split + LOO + repeat를 함께 검증했다.", size=23, bold=True, color=ACCENT)
    labels = ["focus", "tracks", "pairwise", "scenario shift", "threshold tuning"]
    xs = [0.8, 3.0, 5.2, 7.6, 10.0]
    for i, (x, label) in enumerate(zip(xs, labels)):
        add_pipeline_box(slide, x, 1.95, 2.0, 1.0, label, ACCENT if i % 2 == 0 else ACCENT_2)
    add_table(slide, 1.0, 3.45, 5.1, 2.25, [
        ["항목", "값"],
        ["기간", "61일"],
        ["해역", "Houston / NOLA / Seattle"],
        ["own_mmsi", "15"],
        ["pairwise rows", "67,892"],
    ], col_widths=[2.3, 2.4], font_size=14)
    add_callout(slide, 7.2, 3.6, 2.1, 1.0, "Benchmark", "timestamp split", ACCENT)
    add_callout(slide, 9.55, 3.6, 2.1, 1.0, "LOO", "own-ship holdout", ACCENT_2)
    add_callout(slide, 8.38, 4.9, 2.1, 1.0, "Repeat", "case repeat", ACCENT)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "모델 결과", "Results")
    add_textbox(slide, 0.7, 1.1, 11.6, 0.55, "주력 모델은 hgbt, 설명 가능한 비교군은 logreg다.", size=24, bold=True, color=ACCENT)
    add_table(slide, 1.1, 2.0, 7.8, 2.4, [
        ["Model", "Benchmark F1", "ECE", "LOO F1 mean"],
        ["hgbt", "0.9453", "0.0187", "0.9241"],
        ["logreg", "0.8750", "0.0644", "0.8953"],
    ], col_widths=[1.6, 2.0, 1.4, 2.2], font_size=15)
    add_callout(slide, 9.3, 2.2, 2.7, 1.0, "Primary", "hgbt", ACCENT)
    add_callout(slide, 9.3, 3.55, 2.7, 1.0, "Comparator", "logreg", ACCENT_2)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Threshold 결과", "Threshold")
    add_textbox(slide, 0.7, 1.05, 11.7, 0.55, "single best value 대신 shortlist가 더 정직했다.", size=24, bold=True, color=ACCENT)
    add_table(slide, 0.85, 1.95, 4.6, 2.3, [
        ["항목", "값"],
        ["majority profile", "s0p30_w0p55"],
        ["majority ratio", "0.1803"],
        ["mean top-k Jaccard", "0.1388"],
    ], col_widths=[2.2, 2.1], font_size=14)
    add_callout(slide, 5.9, 2.0, 2.0, 1.0, "default", "s0p30_w0p55", ACCENT)
    add_callout(slide, 8.15, 2.0, 2.0, 1.0, "balanced", "s0p30_w0p65", ACCENT_2)
    add_callout(slide, 10.4, 2.0, 2.0, 1.0, "tight", "s0p35_w0p65", ACCENT)
    add_textbox(slide, 0.95, 4.8, 4.2, 0.8, "unstable -> shortlist 운영", size=21, bold=True, color=ACCENT_2)
    slide.shapes.add_picture(str(SCENARIO), Inches(5.7), Inches(3.1), width=Inches(6.1))

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Holdout Figure", "Visual")
    add_textbox(slide, 0.7, 1.0, 11.8, 0.55, "shortlist는 위험 방향보다 contour 민감도를 조정한다.", size=24, bold=True, color=ACCENT)
    slide.shapes.add_picture(str(HOUSTON), Inches(0.7), Inches(1.75), width=Inches(7.1))
    slide.shapes.add_picture(str(NOLA), Inches(8.15), Inches(1.95), width=Inches(2.2))
    slide.shapes.add_picture(str(SEATTLE), Inches(10.5), Inches(1.95), width=Inches(2.2))
    add_bullets(slide, 8.1, 4.65, 4.6, 1.5, [
        "default -> balanced: warning area 축소",
        "balanced -> tight: caution area 추가 축소",
        "dominant sector는 대체로 유지",
    ], size=15)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "기여와 한계", "Close")
    add_textbox(slide, 0.7, 1.0, 11.8, 0.55, "강점은 과장 없는 AIS-only decision-support framing이다.", size=24, bold=True, color=ACCENT)
    add_callout(slide, 0.8, 1.95, 5.7, 2.7, "기여", "spatial field\nscenario-aware contour\nregional failure-mode reporting", ACCENT)
    add_callout(slide, 6.8, 1.95, 5.7, 2.7, "비주장", "autonomy\nlegal safety boundary\nsingle optimum threshold", ACCENT_2)
    add_textbox(slide, 1.0, 5.0, 11.2, 0.8, "Conclusion: explainable spatial decision support", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)


def main():
    ensure_assets()

    prs, blank = base_presentation()
    build_8_slide(prs, blank)
    prs.save(str(OUT_8))

    prs3, blank3 = base_presentation()
    build_3_slide_set(prs3, blank3)
    prs3.save(str(OUT_3))

    print(OUT_8)
    print(OUT_3)


if __name__ == "__main__":
    main()
